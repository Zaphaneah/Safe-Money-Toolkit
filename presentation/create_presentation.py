#!/usr/bin/env python3
"""
Create the 'The Cross That Shows Love' PowerPoint presentation.
Week 2 of the 'All In With Love' Bible class series on Philippians 2:5-8.

Professional design with jewel-tone palette, custom imagery, and clean typography.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ============================================================
# CONSTANTS & COLORS
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

IMAGE_DIR = os.path.join(os.path.dirname(__file__), "images")
OUTPUT_FILE = os.path.join(os.path.dirname(__file__), "The_Cross_That_Shows_Love_Week2.pptx")

# Color palette (jewel tones)
DEEP_PURPLE = RGBColor(0x30, 0x0F, 0x48)
RICH_PURPLE = RGBColor(0x4B, 0x19, 0x6E)
BURGUNDY = RGBColor(0x80, 0x00, 0x20)
DARK_BURGUNDY = RGBColor(0x50, 0x0A, 0x19)
GOLD = RGBColor(0xDA, 0xA5, 0x20)
LIGHT_GOLD = RGBColor(0xFF, 0xD7, 0x64)
PALE_GOLD = RGBColor(0xFF, 0xEB, 0xB4)
DEEP_TEAL = RGBColor(0x00, 0x50, 0x50)
WARM_WHITE = RGBColor(0xFF, 0xF8, 0xF0)
CREAM = RGBColor(0xFF, 0xF5, 0xE6)
NEAR_BLACK = RGBColor(0x0F, 0x0A, 0x14)
DARK_BG = RGBColor(0x14, 0x0C, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
SOFT_GOLD = RGBColor(0xC8, 0x96, 0x2D)


def set_slide_bg(slide, color):
    """Set a solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=24,
                font_color=WARM_WHITE, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Add a text box with formatted text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    # Set vertical alignment
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    # Clear default paragraph and set text
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, default_size=24,
                          default_color=WARM_WHITE, default_bold=False,
                          default_font="Calibri", alignment=PP_ALIGN.LEFT,
                          line_spacing=1.3):
    """
    Add a textbox with multiple formatted lines.
    lines: list of dicts with keys: text, size, color, bold, italic, font, spacing_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            line_info = {"text": line_info}

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line_info.get("text", "")
        p.font.size = Pt(line_info.get("size", default_size))
        p.font.color.rgb = line_info.get("color", default_color)
        p.font.bold = line_info.get("bold", default_bold)
        p.font.italic = line_info.get("italic", False)
        p.font.name = line_info.get("font", default_font)
        p.alignment = line_info.get("alignment", alignment)
        p.space_after = Pt(line_info.get("spacing_after", 4))
        p.space_before = Pt(line_info.get("spacing_before", 0))

    return txBox


def add_shape_rectangle(slide, left, top, width, height, fill_color, border_color=None):
    """Add a filled rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_image(slide, image_name, left, top, width, height):
    """Add an image to the slide."""
    img_path = os.path.join(IMAGE_DIR, image_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
        return True
    else:
        print(f"  WARNING: Image not found: {img_path}")
        return False


def add_line(slide, start_x, start_y, end_x, end_y, color=GOLD, width=1.5):
    """Add a line shape to the slide."""
    from pptx.enum.shapes import MSO_SHAPE
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        start_x, start_y,
        end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


# ============================================================
# SLIDE BUILDERS
# ============================================================

def create_slide1_title(prs):
    """SLIDE 1: Title Slide - elegant bold design on solid background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, DEEP_PURPLE)

    # Top accent bar
    add_shape_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), GOLD)

    # Bottom accent bar
    add_shape_rectangle(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), GOLD)

    # Decorative gold line (upper)
    add_shape_rectangle(slide, Inches(3.5), Inches(1.8), Inches(6.333), Pt(2), GOLD)

    # Series title
    add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10.333), Inches(0.6),
                "ALL IN WITH LOVE  \u2022  WEEK TWO",
                font_size=20, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Calibri",
                line_spacing=1.0)

    # Decorative gold line (below series title)
    add_shape_rectangle(slide, Inches(3.5), Inches(2.0), Inches(6.333), Pt(2), GOLD)

    # Main title
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11.333), Inches(2.0),
                "The Cross That\nShows Love",
                font_size=64, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.15)

    # Scripture reference
    add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.7),
                "Philippians 2:5\u20138",
                font_size=28, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.0)

    # Decorative divider
    add_shape_rectangle(slide, Inches(5.5), Inches(5.6), Inches(2.333), Pt(1.5), GOLD)

    # Church attribution
    add_textbox(slide, Inches(1.5), Inches(6.0), Inches(10.333), Inches(0.5),
                "Pembroke Park Church of Christ",
                font_size=18, font_color=LIGHT_GOLD, italic=False,
                alignment=PP_ALIGN.CENTER, font_name="Calibri",
                line_spacing=1.0)

    # Subtle diamond ornaments
    add_textbox(slide, Inches(5.2), Inches(5.45), Inches(1), Inches(0.4),
                "\u25C6",
                font_size=12, font_color=GOLD, alignment=PP_ALIGN.CENTER,
                line_spacing=1.0)
    add_textbox(slide, Inches(7.1), Inches(5.45), Inches(1), Inches(0.4),
                "\u25C6",
                font_size=12, font_color=GOLD, alignment=PP_ALIGN.CENTER,
                line_spacing=1.0)

    print("  Created: Slide 1 - Title")


def create_slide2_mindset(prs):
    """SLIDE 2: The Mindset of Christ."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image (full bleed, left half)
    add_image(slide, "slide2_cross_mindset.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Semi-transparent overlay on right side for text
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DEEP_PURPLE)

    # Gold accent stripe
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Section number
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "SECTION I",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(1.0),
                "The Mindset of Christ",
                font_size=42, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.2), Inches(3), Pt(2), GOLD)

    # Key quote
    add_textbox(slide, Inches(6.5), Inches(2.6), Inches(6), Inches(1.0),
                "\u201CHave this mind among yourselves,\nwhich is yours in Christ Jesus\u2026\u201D",
                font_size=22, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.3)

    # Scripture ref
    add_textbox(slide, Inches(6.5), Inches(3.7), Inches(6), Inches(0.4),
                "\u2014 Philippians 2:5",
                font_size=16, font_color=MID_GRAY, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Contrast box
    add_multiline_textbox(slide, Inches(6.5), Inches(4.3), Inches(5.5), Inches(1.8),
        [
            {"text": "The Contrast:", "size": 20, "color": GOLD, "bold": True, "spacing_after": 12},
            {"text": "The World says: \u201CClimb UP\u201D", "size": 22, "color": WARM_WHITE, "bold": False, "spacing_after": 8},
            {"text": "Christ says: \u201CDescend DOWN\u201D", "size": 22, "color": LIGHT_GOLD, "bold": True, "spacing_after": 12},
        ],
        alignment=PP_ALIGN.LEFT)

    # Greek term box
    add_shape_rectangle(slide, Inches(6.5), Inches(6.2), Inches(5.5), Inches(0.7), DARK_BURGUNDY)
    add_multiline_textbox(slide, Inches(6.7), Inches(6.25), Inches(5.1), Inches(0.6),
        [
            {"text": "phron\u0113te (\u03c6\u03c1\u03bf\u03bd\u03b5\u1fd6\u03c4\u03b5)  \u2014  \"set your mind on\"",
             "size": 18, "color": PALE_GOLD, "italic": True, "font": "Georgia"},
        ],
        alignment=PP_ALIGN.LEFT)

    print("  Created: Slide 2 - The Mindset of Christ")


def create_slide3_incarnation(prs):
    """SLIDE 3: God With Us - The Incarnation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image (full width, behind)
    add_image(slide, "slide3_incarnation.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DEEP_PURPLE)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Section number
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "SECTION II",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_multiline_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(1.2),
        [
            {"text": "God With Us\u2014", "size": 40, "color": WARM_WHITE, "bold": True, "font": "Georgia", "spacing_after": 2},
            {"text": "The Incarnation", "size": 40, "color": LIGHT_GOLD, "bold": True, "font": "Georgia"},
        ],
        alignment=PP_ALIGN.LEFT)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.5), Inches(3), Pt(2), GOLD)

    # Three key points
    add_multiline_textbox(slide, Inches(6.5), Inches(2.9), Inches(5.8), Inches(3.5),
        [
            {"text": "\u201CForm of God\u201D (morph\u0113)", "size": 24, "color": LIGHT_GOLD, "bold": True, "spacing_after": 4, "font": "Georgia"},
            {"text": "Essential divine nature (v. 6)", "size": 19, "color": CREAM, "spacing_after": 16},
            {"text": "Did not grasp equality with God", "size": 24, "color": LIGHT_GOLD, "bold": True, "spacing_after": 4, "font": "Georgia"},
            {"text": "Refused to cling to divine privileges", "size": 19, "color": CREAM, "spacing_after": 16},
            {"text": "The Infinite became Finite", "size": 24, "color": LIGHT_GOLD, "bold": True, "spacing_after": 4, "font": "Georgia"},
            {"text": "Complete identification with humanity", "size": 19, "color": CREAM, "spacing_after": 8},
        ],
        alignment=PP_ALIGN.LEFT)

    # Emphasis box
    add_shape_rectangle(slide, Inches(6.5), Inches(5.8), Inches(5.5), Inches(0.9), DARK_BURGUNDY)
    add_textbox(slide, Inches(6.7), Inches(5.9), Inches(5.1), Inches(0.7),
                "The Creator entered His own creation\u2014\ncomplete identification with humanity",
                font_size=18, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.3)

    print("  Created: Slide 3 - God With Us")


def create_slide4_emptied(prs):
    """SLIDE 4: He Emptied Himself - Kenosis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide4_emptying.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, RICH_PURPLE)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Section number
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "SECTION III",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.8),
                "He Emptied Himself",
                font_size=42, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.1), Inches(3), Pt(2), GOLD)

    # Keyword highlight
    add_shape_rectangle(slide, Inches(6.5), Inches(2.4), Inches(5.5), Inches(0.8), DARK_BURGUNDY)
    add_multiline_textbox(slide, Inches(6.7), Inches(2.45), Inches(5.1), Inches(0.7),
        [
            {"text": "KEN\u014cSIS (\u03ba\u03ad\u03bd\u03c9\u03c3\u03b9\u03c2)", "size": 28, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 2},
            {"text": "Greek: self-emptying", "size": 16, "color": CREAM, "italic": True},
        ],
        alignment=PP_ALIGN.LEFT)

    # Three descending points with arrow visual
    add_multiline_textbox(slide, Inches(6.5), Inches(3.5), Inches(5.8), Inches(2.5),
        [
            {"text": "\u25BC  Gave up privileges of deity", "size": 22, "color": WARM_WHITE, "bold": False, "spacing_after": 14, "font": "Calibri"},
            {"text": "\u25BC  Took form of a servant (doulos)", "size": 22, "color": WARM_WHITE, "bold": False, "spacing_after": 14, "font": "Calibri"},
            {"text": "\u25BC  Born in likeness of men", "size": 22, "color": WARM_WHITE, "bold": False, "spacing_after": 14, "font": "Calibri"},
        ],
        alignment=PP_ALIGN.LEFT)

    # Key phrase
    add_shape_rectangle(slide, Inches(6.5), Inches(5.8), Inches(5.5), Inches(0.9), DEEP_PURPLE)
    add_textbox(slide, Inches(6.7), Inches(5.9), Inches(5.1), Inches(0.7),
                "\u201CNot emptying deity,\nbut the privileges of deity\u201D",
                font_size=20, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.3)

    print("  Created: Slide 4 - He Emptied Himself")


def create_slide5_humbled(prs):
    """SLIDE 5: He Humbled Himself."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide5_humble.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DARK_BURGUNDY)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Section number
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "SECTION IV",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.8),
                "He Humbled Himself",
                font_size=42, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.1), Inches(3), Pt(2), GOLD)

    # Emphasis subtitle
    add_textbox(slide, Inches(6.5), Inches(2.4), Inches(6), Inches(0.6),
                "Active Humility, Not Passive Humiliation",
                font_size=22, font_color=LIGHT_GOLD, bold=True, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Key points
    add_multiline_textbox(slide, Inches(6.5), Inches(3.2), Inches(5.8), Inches(2.5),
        [
            {"text": "\u201CHumbled himself\u201D", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 4},
            {"text": "He CHOSE this path (v. 8a)", "size": 20, "color": CREAM, "spacing_after": 18},
            {"text": "Found in human form", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 4},
            {"text": "No special recognition or status", "size": 20, "color": CREAM, "spacing_after": 18},
        ],
        alignment=PP_ALIGN.LEFT)

    # Isaiah quote box
    add_shape_rectangle(slide, Inches(6.5), Inches(5.5), Inches(5.5), Inches(1.2), DEEP_PURPLE)
    add_multiline_textbox(slide, Inches(6.7), Inches(5.6), Inches(5.1), Inches(1.0),
        [
            {"text": "\u201CHe had no form or majesty that we should\nlook at him, and no beauty that we should\ndesire him.\u201D",
             "size": 18, "color": PALE_GOLD, "italic": True, "font": "Georgia", "spacing_after": 4},
            {"text": "\u2014 Isaiah 53:2", "size": 14, "color": MID_GRAY, "italic": True},
        ],
        alignment=PP_ALIGN.LEFT)

    print("  Created: Slide 5 - He Humbled Himself")


def create_slide6_death(prs):
    """SLIDE 6: Obedient to Death."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide6_cross_death.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, NEAR_BLACK)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, BURGUNDY)

    # Section number
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "SECTION V",
                font_size=14, font_color=BURGUNDY, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.8),
                "Obedient to Death",
                font_size=42, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.1), Inches(3), Pt(2), BURGUNDY)

    # Central quote
    add_shape_rectangle(slide, Inches(6.5), Inches(2.4), Inches(5.5), Inches(0.8), DARK_BURGUNDY)
    add_textbox(slide, Inches(6.7), Inches(2.5), Inches(5.1), Inches(0.6),
                "\u201CEven death on a cross\u201D  \u2014  v. 8b",
                font_size=24, font_color=PALE_GOLD, italic=True, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.0)

    # Three characteristics
    add_multiline_textbox(slide, Inches(6.5), Inches(3.5), Inches(5.8), Inches(2.5),
        [
            {"text": "Most Shameful Death", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 4},
            {"text": "Reserved for slaves and criminals", "size": 18, "color": CREAM, "spacing_after": 16},
            {"text": "Most Painful Death", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 4},
            {"text": "Designed for maximum suffering", "size": 18, "color": CREAM, "spacing_after": 16},
            {"text": "Cursed by God", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia", "spacing_after": 4},
            {"text": "Deuteronomy 21:23", "size": 18, "color": CREAM, "italic": True, "spacing_after": 8},
        ],
        alignment=PP_ALIGN.LEFT)

    # Bottom emphasis
    add_shape_rectangle(slide, Inches(6.5), Inches(6.2), Inches(5.5), Inches(0.7), BURGUNDY)
    add_textbox(slide, Inches(6.7), Inches(6.3), Inches(5.1), Inches(0.5),
                "Love\u2019s descent to the absolute bottom",
                font_size=22, font_color=WARM_WHITE, bold=True, italic=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.0)

    print("  Created: Slide 6 - Obedient to Death")


def create_slide7_five_movements(prs):
    """SLIDE 7: The Five Descending Movements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide7_descending.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DEEP_TEAL)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(0.6), Inches(6), Inches(0.5),
                "SUMMARY",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    add_textbox(slide, Inches(6.5), Inches(0.95), Inches(6), Inches(0.8),
                "The Five Descending\nMovements",
                font_size=36, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.1)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.2), Inches(3), Pt(2), GOLD)

    # Five numbered items with descending visual
    movements = [
        ("1", "The Mindset", "Chose to think differently"),
        ("2", "God With Us", "The Incarnation"),
        ("3", "Self-Emptying", "Became a servant"),
        ("4", "Humbling", "Chose the lower place"),
        ("5", "Obedience unto Death", "Even the cross"),
    ]

    y_start = 2.6
    for i, (num, title, subtitle) in enumerate(movements):
        y = y_start + i * 0.9
        indent = i * 0.15  # Slight indent for each step

        # Number circle
        circle_left = Inches(6.5 + indent)
        circle_top = Inches(y)
        # Use a shape for the circle
        from pptx.enum.shapes import MSO_SHAPE
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, Inches(0.45), Inches(0.45))
        circle.fill.solid()
        t_ratio = i / 4
        r = int(218 - t_ratio * 90)
        g = int(165 - t_ratio * 120)
        b = int(32 + t_ratio * 30)
        circle.fill.fore_color.rgb = RGBColor(r, g, b)
        circle.line.fill.background()

        # Number text
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NEAR_BLACK
        p.font.name = "Georgia"
        p.alignment = PP_ALIGN.CENTER

        # Title and subtitle
        add_multiline_textbox(slide, Inches(7.2 + indent), Inches(y - 0.05), Inches(5), Inches(0.5),
            [
                {"text": title, "size": 22, "color": WARM_WHITE, "bold": True, "font": "Georgia", "spacing_after": 0},
                {"text": subtitle, "size": 16, "color": CREAM, "italic": True, "spacing_after": 0},
            ],
            alignment=PP_ALIGN.LEFT)

        # Downward arrow between items (except last)
        if i < len(movements) - 1:
            arrow_x = Inches(6.72 + indent + 0.075)
            add_textbox(slide, arrow_x, Inches(y + 0.45), Inches(0.5), Inches(0.3),
                        "\u25BC",
                        font_size=12, font_color=GOLD, alignment=PP_ALIGN.CENTER,
                        line_spacing=1.0)

    print("  Created: Slide 7 - Five Descending Movements")


def create_slide8_pattern(prs):
    """SLIDE 8: The Pattern for Our Lives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide8_serving_hands.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DARK_BURGUNDY)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Label
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "APPLICATION",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.8),
                "The Pattern for\nOur Lives",
                font_size=40, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.1)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.4), Inches(3), Pt(2), GOLD)

    # Key statement
    add_textbox(slide, Inches(6.5), Inches(2.8), Inches(5.5), Inches(1.2),
                "You cannot be \u201Call in\u201D with love while holding back from the downward path.",
                font_size=22, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.4)

    # Emphasis
    add_shape_rectangle(slide, Inches(6.5), Inches(4.0), Inches(5.5), Inches(0.6), BURGUNDY)
    add_textbox(slide, Inches(6.7), Inches(4.05), Inches(5.1), Inches(0.5),
                "The cross is our pattern for living",
                font_size=22, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.0)

    # Three action steps
    steps = [
        ("\u2023  Serve others", "Follow Christ\u2019s example of humble service"),
        ("\u2023  Give up rights", "Release what you\u2019re entitled to for others\u2019 sake"),
        ("\u2023  Choose humility", "Actively take the lower place"),
    ]

    y = 4.9
    for title, desc in steps:
        add_multiline_textbox(slide, Inches(6.5), Inches(y), Inches(5.5), Inches(0.6),
            [
                {"text": title, "size": 22, "color": LIGHT_GOLD, "bold": True, "font": "Calibri", "spacing_after": 2},
                {"text": desc, "size": 16, "color": CREAM, "italic": True, "spacing_after": 4},
            ],
            alignment=PP_ALIGN.LEFT)
        y += 0.75

    print("  Created: Slide 8 - The Pattern for Our Lives")


def create_slide9_challenge(prs):
    """SLIDE 9: This Week's Challenge."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEAR_BLACK)

    # Background image on left
    add_image(slide, "slide9_contemplating.jpg",
              Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT)

    # Text panel on right
    add_shape_rectangle(slide, Inches(5.8), Inches(0), Inches(7.533), SLIDE_HEIGHT, DEEP_PURPLE)
    add_shape_rectangle(slide, Inches(6.0), Inches(0), Pt(4), SLIDE_HEIGHT, GOLD)

    # Label
    add_textbox(slide, Inches(6.5), Inches(0.8), Inches(6), Inches(0.5),
                "CHALLENGE",
                font_size=14, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=1.0)

    # Main heading
    add_textbox(slide, Inches(6.5), Inches(1.2), Inches(6), Inches(0.8),
                "This Week\u2019s\nChallenge",
                font_size=40, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.1)

    # Gold divider
    add_shape_rectangle(slide, Inches(6.5), Inches(2.4), Inches(3), Pt(2), GOLD)

    # Challenge text
    add_shape_rectangle(slide, Inches(6.5), Inches(2.9), Inches(5.5), Inches(1.2), DARK_BURGUNDY)
    add_multiline_textbox(slide, Inches(6.7), Inches(3.0), Inches(5.1), Inches(1.0),
        [
            {"text": "IDENTIFY", "size": 16, "color": GOLD, "bold": True, "spacing_after": 4},
            {"text": "one area where you\u2019re clinging\nto privilege or status", "size": 22, "color": WARM_WHITE, "font": "Georgia", "spacing_after": 0},
        ],
        alignment=PP_ALIGN.LEFT)

    # Action
    add_shape_rectangle(slide, Inches(6.5), Inches(4.4), Inches(5.5), Inches(1.0), RICH_PURPLE)
    add_multiline_textbox(slide, Inches(6.7), Inches(4.5), Inches(5.1), Inches(0.8),
        [
            {"text": "THEN", "size": 16, "color": GOLD, "bold": True, "spacing_after": 4},
            {"text": "Take one step DOWN this week", "size": 24, "color": LIGHT_GOLD, "bold": True, "font": "Georgia"},
        ],
        alignment=PP_ALIGN.LEFT)

    # Question
    add_textbox(slide, Inches(6.5), Inches(5.8), Inches(5.5), Inches(1.0),
                "Where are you\nholding back?",
                font_size=32, font_color=WARM_WHITE, bold=True, italic=True,
                alignment=PP_ALIGN.LEFT, font_name="Georgia",
                line_spacing=1.2)

    print("  Created: Slide 9 - This Week's Challenge")


def create_slide10_closing(prs):
    """SLIDE 10: Closing - powerful text only."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_PURPLE)

    # Top accent bar
    add_shape_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), GOLD)
    # Bottom accent bar
    add_shape_rectangle(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06), GOLD)

    # "The world says" block
    add_textbox(slide, Inches(1.5), Inches(0.7), Inches(10.333), Inches(0.4),
                "THE WORLD SAYS:",
                font_size=18, font_color=MID_GRAY, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Calibri",
                line_spacing=1.0)

    add_multiline_textbox(slide, Inches(1.5), Inches(1.1), Inches(10.333), Inches(1.2),
        [
            {"text": "Look out for yourself.  Climb the ladder.  Protect your rights.",
             "size": 24, "color": LIGHT_GRAY, "font": "Georgia", "italic": True,
             "alignment": PP_ALIGN.CENTER},
        ],
        alignment=PP_ALIGN.CENTER)

    # Visual divider 1
    add_shape_rectangle(slide, Inches(5.5), Inches(2.3), Inches(2.333), Pt(2), GOLD)
    add_textbox(slide, Inches(5.8), Inches(2.15), Inches(1.733), Inches(0.4),
                "\u25C6", font_size=14, font_color=GOLD, alignment=PP_ALIGN.CENTER,
                line_spacing=1.0)

    # "The cross says" block
    add_textbox(slide, Inches(1.5), Inches(2.7), Inches(10.333), Inches(0.4),
                "THE CROSS SAYS:",
                font_size=18, font_color=GOLD, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Calibri",
                line_spacing=1.0)

    add_multiline_textbox(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(1.2),
        [
            {"text": "Give yourself away.  Descend to serve.  Renounce your rights.",
             "size": 26, "color": WARM_WHITE, "font": "Georgia", "italic": True, "bold": True,
             "alignment": PP_ALIGN.CENTER},
        ],
        alignment=PP_ALIGN.CENTER)

    # Visual divider 2
    add_shape_rectangle(slide, Inches(5.5), Inches(4.3), Inches(2.333), Pt(2), GOLD)
    add_textbox(slide, Inches(5.8), Inches(4.15), Inches(1.733), Inches(0.4),
                "\u25C6", font_size=14, font_color=GOLD, alignment=PP_ALIGN.CENTER,
                line_spacing=1.0)

    # Contrast statement
    add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10.333), Inches(0.8),
                "These are opposite paths.\nYou cannot walk both.",
                font_size=28, font_color=PALE_GOLD, italic=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.3)

    # Final call - prominent
    add_shape_rectangle(slide, Inches(3), Inches(5.8), Inches(7.333), Inches(1.0), BURGUNDY)
    add_textbox(slide, Inches(3.2), Inches(5.9), Inches(6.933), Inches(0.8),
                "Choose the way of the cross.",
                font_size=36, font_color=WARM_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Georgia",
                line_spacing=1.0)

    print("  Created: Slide 10 - Closing")


# ============================================================
# MAIN
# ============================================================

def main():
    print("Creating presentation: The Cross That Shows Love")
    print("=" * 55)

    prs = Presentation()

    # Set 16:9 widescreen
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create all slides
    create_slide1_title(prs)
    create_slide2_mindset(prs)
    create_slide3_incarnation(prs)
    create_slide4_emptied(prs)
    create_slide5_humbled(prs)
    create_slide6_death(prs)
    create_slide7_five_movements(prs)
    create_slide8_pattern(prs)
    create_slide9_challenge(prs)
    create_slide10_closing(prs)

    # Save presentation
    prs.save(OUTPUT_FILE)
    file_size = os.path.getsize(OUTPUT_FILE)
    print(f"\n{'=' * 55}")
    print(f"Presentation saved: {OUTPUT_FILE}")
    print(f"File size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
    print(f"Slides: 10")
    print(f"Format: .pptx (16:9 widescreen)")


if __name__ == "__main__":
    main()
